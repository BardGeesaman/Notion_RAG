"""PowerPoint slide generator for experiments and datasets."""
from __future__ import annotations

from typing import List
from uuid import UUID
from io import BytesIO

from pptx import Presentation
from pptx.util import Pt

from amprenta_rag.logging_utils import get_logger

logger = get_logger(__name__)


def add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    """
    Add a title slide to the presentation.
    
    Args:
        prs: Presentation object
        title: Main title text
        subtitle: Optional subtitle text
    """
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    if subtitle:
        subtitle_shape.text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: List[str]) -> None:
    """
    Add a bullet point slide to the presentation.
    
    Args:
        prs: Presentation object
        title: Slide title
        bullets: List of bullet point strings
    """
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    for bullet_text in bullets:
        p = text_frame.add_paragraph()
        p.text = bullet_text
        p.level = 0
        p.font.size = Pt(18)


def generate_experiment_slides(experiment_id: UUID, db) -> bytes:
    """
    Generate PowerPoint slides for an experiment.
    
    Args:
        experiment_id: UUID of the experiment
        db: Database session
        
    Returns:
        Bytes of the PowerPoint presentation
    """
    from amprenta_rag.database.models import Experiment
    
    # Load experiment from DB
    experiment = db.query(Experiment).filter(Experiment.id == experiment_id).first()
    if not experiment:
        raise ValueError(f"Experiment {experiment_id} not found")
    
    # Create presentation
    prs = Presentation()
    
    # Title slide
    experiment_name = experiment.name or "Unknown Experiment"
    date_str = experiment.created_at.strftime("%Y-%m-%d") if experiment.created_at else "Unknown Date"
    add_title_slide(prs, experiment_name, f"Created: {date_str}")
    
    # Overview slide
    overview_bullets = []
    if experiment.description:
        overview_bullets.append(f"Description: {experiment.description}")
    if experiment.design_type:
        overview_bullets.append(f"Design Type: {experiment.design_type}")
    
    # Get organism from datasets if available
    organism_info = "Unknown"
    if hasattr(experiment, "datasets") and experiment.datasets:
        organisms = set()
        for dataset in experiment.datasets:
            if hasattr(dataset, "organism") and dataset.organism:
                if isinstance(dataset.organism, list):
                    organisms.update(dataset.organism)
                else:
                    organisms.add(dataset.organism)
        if organisms:
            organism_info = ", ".join(sorted(organisms))
    
    overview_bullets.append(f"Organism: {organism_info}")
    
    if overview_bullets:
        add_bullet_slide(prs, "Overview", overview_bullets)
    
    # Summary slide
    dataset_count = len(experiment.datasets) if hasattr(experiment, "datasets") else 0
    summary_bullets = [
        f"Total Datasets: {dataset_count}",
    ]
    
    if experiment.disease:
        disease_str = ", ".join(experiment.disease) if isinstance(experiment.disease, list) else str(experiment.disease)
        summary_bullets.append(f"Disease: {disease_str}")
    
    if experiment.matrix:
        matrix_str = ", ".join(experiment.matrix) if isinstance(experiment.matrix, list) else str(experiment.matrix)
        summary_bullets.append(f"Matrix: {matrix_str}")
    
    if experiment.sample_groups:
        summary_bullets.append("Sample Groups: See design metadata")
    
    add_bullet_slide(prs, "Summary", summary_bullets)
    
    # Save to BytesIO
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output.read()


def generate_dataset_slides(dataset_id: UUID, db) -> bytes:
    """
    Generate PowerPoint slides for a dataset.
    
    Args:
        dataset_id: UUID of the dataset
        db: Database session
        
    Returns:
        Bytes of the PowerPoint presentation
    """
    from amprenta_rag.database.models import Dataset
    
    # Load dataset from DB
    dataset = db.query(Dataset).filter(Dataset.id == dataset_id).first()
    if not dataset:
        raise ValueError(f"Dataset {dataset_id} not found")
    
    # Create presentation
    prs = Presentation()
    
    # Title slide
    dataset_name = dataset.name or "Unknown Dataset"
    date_str = dataset.created_at.strftime("%Y-%m-%d") if dataset.created_at else "Unknown Date"
    add_title_slide(prs, dataset_name, f"Created: {date_str}")
    
    # Overview slide
    overview_bullets = []
    
    if dataset.omics_type:
        overview_bullets.append(f"Omics Type: {dataset.omics_type}")
    
    if dataset.dataset_source_type:
        overview_bullets.append(f"Source Type: {dataset.dataset_source_type}")
    elif dataset.data_origin:
        overview_bullets.append(f"Data Origin: {dataset.data_origin}")
    
    if dataset.description:
        overview_bullets.append(f"Description: {dataset.description}")
    elif dataset.summary:
        overview_bullets.append(f"Summary: {dataset.summary}")
    
    if dataset.organism:
        organism_str = ", ".join(dataset.organism) if isinstance(dataset.organism, list) else str(dataset.organism)
        overview_bullets.append(f"Organism: {organism_str}")
    
    if overview_bullets:
        add_bullet_slide(prs, "Overview", overview_bullets)
    
    # Additional details slide
    details_bullets = []
    
    if dataset.disease:
        disease_str = ", ".join(dataset.disease) if isinstance(dataset.disease, list) else str(dataset.disease)
        details_bullets.append(f"Disease: {disease_str}")
    
    if dataset.sample_type:
        sample_type_str = ", ".join(dataset.sample_type) if isinstance(dataset.sample_type, list) else str(dataset.sample_type)
        details_bullets.append(f"Sample Type: {sample_type_str}")
    
    if dataset.sample_group:
        details_bullets.append(f"Sample Group: {dataset.sample_group}")
    
    if dataset.timepoint:
        details_bullets.append(f"Timepoint: {dataset.timepoint}")
    
    if dataset.intervention:
        details_bullets.append(f"Intervention: {dataset.intervention}")
    
    if details_bullets:
        add_bullet_slide(prs, "Additional Details", details_bullets)
    
    # Save to BytesIO
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output.read()

