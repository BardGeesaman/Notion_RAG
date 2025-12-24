"""Parsers for extracting structured text/metadata from common office file types."""

from __future__ import annotations

from pathlib import Path
from typing import Any, Callable, Dict, List

import pandas as pd


def parse_docx(file_path: str) -> dict:
    """
    Parse a .docx file using python-docx.

    Returns:
        {"text": str, "tables": list, "metadata": dict}
    """
    try:
        from docx import Document  # type: ignore
    except ImportError as e:  # pragma: no cover
        raise ImportError("python-docx is required to parse .docx files") from e

    doc = Document(file_path)

    paragraphs = [p.text for p in doc.paragraphs if (p.text or "").strip()]
    text = "\n".join(paragraphs)

    tables: List[Dict[str, Any]] = []
    for t in doc.tables:
        rows: List[List[str]] = []
        for row in t.rows:
            rows.append([cell.text for cell in row.cells])
        tables.append(
            {
                "n_rows": len(rows),
                "n_cols": max((len(r) for r in rows), default=0),
                "rows": rows,
            }
        )

    cp = doc.core_properties
    metadata = {
        "author": getattr(cp, "author", None),
        "title": getattr(cp, "title", None),
        "created": getattr(cp, "created", None),
        "modified": getattr(cp, "modified", None),
    }
    return {"text": text, "tables": tables, "metadata": metadata}


def parse_pptx(file_path: str) -> dict:
    """
    Parse a .pptx file using python-pptx.

    Returns:
        {"slides": [{slide_num, title, content, notes}], "metadata": dict}
    """
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as e:  # pragma: no cover
        raise ImportError("python-pptx is required to parse .pptx files") from e

    prs = Presentation(file_path)

    slides: List[Dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = None
        if getattr(slide.shapes, "title", None) is not None:
            try:
                title = slide.shapes.title.text
            except Exception:
                title = None

        parts: List[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "has_text_frame"):
                continue
            if not shape.has_text_frame:
                continue
            txt = getattr(shape.text_frame, "text", "") or ""
            txt = txt.strip()
            if txt:
                parts.append(txt)

        notes_text = ""
        try:
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            notes_text = ""

        slides.append(
            {
                "slide_num": idx,
                "title": title,
                "content": "\n".join(parts),
                "notes": notes_text,
            }
        )

    cp = prs.core_properties
    metadata = {
        "author": getattr(cp, "author", None),
        "title": getattr(cp, "title", None),
        "created": getattr(cp, "created", None),
        "modified": getattr(cp, "modified", None),
    }
    return {"slides": slides, "metadata": metadata}


def parse_excel(file_path: str) -> dict:
    """
    Parse an .xlsx file using pandas + openpyxl.

    Returns:
        {"sheets": [{name, columns, row_count, sample_rows}], "metadata": dict}
    """
    path = Path(file_path)
    xls = pd.ExcelFile(path, engine="openpyxl")

    sheets: List[Dict[str, Any]] = []
    for name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=name)
        sheets.append(
            {
                "name": name,
                "columns": [str(c) for c in df.columns.tolist()],
                "row_count": int(df.shape[0]),
                "sample_rows": df.head(10).to_dict(orient="records"),
            }
        )

    metadata = {"sheet_count": len(xls.sheet_names), "sheet_names": list(xls.sheet_names)}
    return {"sheets": sheets, "metadata": metadata}


def parse_csv(file_path: str) -> dict:
    """
    Parse a .csv file using pandas.

    Returns:
        {"columns": list, "row_count": int, "sample": list}
    """
    df = pd.read_csv(file_path)
    return {
        "columns": [str(c) for c in df.columns.tolist()],
        "row_count": int(df.shape[0]),
        "sample": df.head(10).to_dict(orient="records"),
    }


def get_parser(file_path: str) -> Callable[[str], dict]:
    """
    Auto-detect a parser based on file extension.

    Raises:
        ValueError for unsupported types.
    """
    suffix = Path(file_path).suffix.lower().lstrip(".")
    if suffix == "docx":
        return parse_docx
    if suffix == "pptx":
        return parse_pptx
    if suffix in {"xlsx", "xls"}:
        return parse_excel
    if suffix == "csv":
        return parse_csv
    raise ValueError(f"Unsupported file type: .{suffix}")


